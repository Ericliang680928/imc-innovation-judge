"""
檔案讀取模組:支援 PDF / DOCX / PPTX / TXT / MD
"""

import io
from pathlib import Path


def read_pdf(file_bytes: bytes) -> str:
    """讀取 PDF,優先用 pdfplumber,失敗則用 pypdf。"""
    text_parts = []
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    text_parts.append(f"\n----- 第 {i} 頁 -----\n{page_text}")
        full_text = "\n".join(text_parts).strip()
        if full_text:
            return full_text
    except Exception as e:
        text_parts.append(f"[pdfplumber 失敗:{e}]")

    # fallback: pypdf
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(file_bytes))
        parts = []
        for i, page in enumerate(reader.pages, 1):
            page_text = page.extract_text() or ""
            if page_text.strip():
                parts.append(f"\n----- 第 {i} 頁 -----\n{page_text}")
        return "\n".join(parts).strip()
    except Exception as e:
        raise RuntimeError(f"PDF 讀取失敗:{e}")


def read_docx(file_bytes: bytes) -> str:
    """讀取 Word .docx,擷取段落與表格內容。"""
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    parts = []

    # 段落
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    # 表格
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip(" |"):
                parts.append(row_text)

    return "\n".join(parts).strip()


def read_pptx(file_bytes: bytes) -> str:
    """讀取 PowerPoint .pptx,逐張投影片擷取文字與備註。"""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []

    for i, slide in enumerate(prs.slides, 1):
        slide_parts = [f"\n===== 投影片 {i} ====="]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            slide_parts.append(run.text)
            # 表格
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        slide_parts.append(row_text)
        # 講者備註
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"[講者備註] {notes}")
        parts.append("\n".join(slide_parts))

    return "\n".join(parts).strip()


def read_text(file_bytes: bytes) -> str:
    """讀取 .txt / .md(嘗試多種編碼)。"""
    for enc in ("utf-8", "utf-8-sig", "big5", "cp950", "gb18030"):
        try:
            return file_bytes.decode(enc).strip()
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="replace").strip()


def read_file(file_name: str, file_bytes: bytes) -> str:
    """根據副檔名自動分派到對應 reader。"""
    ext = Path(file_name).suffix.lower()
    if ext == ".pdf":
        return read_pdf(file_bytes)
    elif ext == ".docx":
        return read_docx(file_bytes)
    elif ext == ".pptx":
        return read_pptx(file_bytes)
    elif ext in (".txt", ".md"):
        return read_text(file_bytes)
    else:
        raise ValueError(f"不支援的副檔名:{ext}(支援 .pdf / .docx / .pptx / .txt / .md)")
